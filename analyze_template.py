"""Analyze PowerPoint template structure."""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt


def analyze_template(pptx_path: str) -> None:
    """Analyze the structure of a PowerPoint template."""
    prs = Presentation(pptx_path)

    print("=" * 60)
    print(f"TEMPLATE ANALYSIS: {pptx_path}")
    print("=" * 60)

    # Slide dimensions
    print(f"\nSlide dimensions: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")

    # Slide layouts
    print(f"\n{'=' * 40}")
    print("SLIDE LAYOUTS AVAILABLE:")
    print("=" * 40)
    for i, layout in enumerate(prs.slide_layouts):
        print(f"\n[{i}] {layout.name}")
        for ph in layout.placeholders:
            print(f"    - Placeholder {ph.placeholder_format.idx}: {ph.placeholder_format.type} ({ph.name})")

    # Existing slides
    print(f"\n{'=' * 40}")
    print(f"EXISTING SLIDES: {len(prs.slides)}")
    print("=" * 40)

    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i + 1} ---")
        print(f"Layout: {slide.slide_layout.name}")

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text[:50].replace("\n", " ")
                print(f"  [{shape.shape_type}] {shape.name}: \"{text}...\"" if len(shape.text_frame.text) > 50 else f"  [{shape.shape_type}] {shape.name}: \"{text}\"")
            else:
                print(f"  [{shape.shape_type}] {shape.name}")


if __name__ == "__main__":
    analyze_template("/Users/jainal09/Downloads/MS_SA_Template_Final.pptx")
