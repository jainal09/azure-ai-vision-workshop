"""
Generate Azure AI Vision Workshop PowerPoint from template.

This script takes the Microsoft Student Ambassadors branded template
and populates it with the workshop content with proper formatting.
"""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


# Template and output paths
TEMPLATE_PATH = "/Users/jainal09/Downloads/MS_SA_Template_Final.pptx"
OUTPUT_PATH = "/Users/jainal09/msa-q4-25/Azure_AI_Vision_Workshop.pptx"


def clear_slides(prs: Presentation, keep_first_n: int = 0) -> None:
    """Remove all slides except the first n."""
    slide_ids = [slide.slide_id for slide in prs.slides]
    for i, slide_id in enumerate(slide_ids):
        if i >= keep_first_n:
            rId = prs.slides._sldIdLst[keep_first_n].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[keep_first_n]


def get_layout_index(layout_type: str) -> int:
    """Map layout type to template layout index."""
    layout_map = {
        "title": 3,  # Walkin 4 - title slide
        "content": 6,  # Title & Non-bulleted text
        "content_alt": 7,  # 2_Title & Non-bulleted text
        "content_alt2": 8,  # 3_Title & Non-bulleted text
        "bullets": 9,  # Title and Content (with bullets)
        "two_column": 11,  # Two Column Non-bulleted text
        "demo": 19,  # Demo slide
        "video": 20,  # Video slide
        "section": 21,  # Section Title
        "code": 27,  # Developer Code Layout
        "closing": 28,  # Closing logo slide_1
    }
    return layout_map.get(layout_type, 6)


def format_text_frame(text_frame, content: str, is_bullet_list: bool = False) -> None:
    """Format text frame with proper styling."""
    text_frame.clear()

    lines = content.strip().split("\n")

    for i, line in enumerate(lines):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        # Detect line type and format accordingly
        stripped = line.strip()

        if not stripped:
            # Empty line - add spacing
            p.text = ""
            p.space_after = Pt(12)
            continue

        # Check if it's a header line (no bullet, followed by content)
        is_header = (
            stripped and
            not stripped.startswith("-") and
            not stripped.startswith("•") and
            not stripped.startswith("·") and
            len(stripped) < 50 and
            ":" not in stripped and
            i + 1 < len(lines) and
            (lines[i + 1].strip().startswith("-") or lines[i + 1].strip().startswith("•") or not lines[i + 1].strip())
        )

        # Check if it's a section header (ends with :)
        is_section_header = stripped.endswith(":") and len(stripped) < 60

        # Check if it's a bullet point
        is_bullet = stripped.startswith("-") or stripped.startswith("•") or stripped.startswith("·")

        if is_bullet:
            # Bullet point
            p.text = "• " + stripped.lstrip("-•· ")
            p.level = 0
            p.font.size = Pt(18)
            p.space_before = Pt(6)
            p.space_after = Pt(6)
        elif is_section_header:
            # Section header (bold, slightly larger)
            p.text = stripped
            p.font.bold = True
            p.font.size = Pt(20)
            p.space_before = Pt(18)
            p.space_after = Pt(6)
        elif is_header:
            # Content header
            p.text = stripped
            p.font.bold = True
            p.font.size = Pt(22)
            p.space_before = Pt(14)
            p.space_after = Pt(4)
        else:
            # Regular text
            p.text = stripped
            p.font.size = Pt(18)
            p.space_before = Pt(4)
            p.space_after = Pt(4)


def add_slide(prs: Presentation, slide_data: dict) -> None:
    """Add a slide with the specified layout and content."""
    layout_idx = get_layout_index(slide_data["layout"])
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)

    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = slide_data.get("title", "")
        # Make title bold
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.bold = True

    # Set content based on layout type
    if slide_data["layout"] in ["content", "content_alt", "content_alt2", "bullets", "code"]:
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 2:  # BODY type
                content = slide_data.get("content", "")
                format_text_frame(shape.text_frame, content)
                break

    elif slide_data["layout"] == "demo":
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 2:  # BODY type
                shape.text = slide_data.get("subtitle", "")
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(24)
                break

    elif slide_data["layout"] == "title":
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 2:  # BODY type
                shape.text = slide_data.get("subtitle", "")
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(24)
                break


def generate_presentation() -> None:
    """Generate the workshop presentation from template."""
    print(f"Loading template: {TEMPLATE_PATH}")

    # Copy template to output
    shutil.copy(TEMPLATE_PATH, OUTPUT_PATH)

    # Open the copy
    prs = Presentation(OUTPUT_PATH)

    print(f"Template has {len(prs.slides)} existing slides")

    # Clear all slides except the first walk-in slide
    clear_slides(prs, keep_first_n=1)
    print(f"Cleared slides, keeping first 1. Now have {len(prs.slides)} slides")

    # Workshop slides content
    slides_content = [
        {
            "layout": "title",
            "title": "Build an AI-Powered Image Analyzer",
            "subtitle": "Azure AI Vision Workshop",
        },
        {
            "layout": "content",
            "title": "About Me",
            "content": """Jainal Gosaliya

Beta Microsoft Learn Student Ambassador

Pursuing Masters of Software Engineering Systems
at Northeastern University

Passionate about:
- Python
- Artificial Intelligence
- Microsoft Azure

linkedin.com/in/jainal-gosaliya
github.com/jainal09""",
        },
        {
            "layout": "content_alt",
            "title": "Agenda",
            "content": """What We'll Cover Today:

- Introduction to Azure AI Vision
- Key Features & Capabilities
- Real-World Use Cases
- Live Demo: Build an Image Analyzer
- Resources & Next Steps
- Q&A

By the end, you'll be able to:
- Set up Azure AI Vision resource from scratch
- Build a working image analyzer application
- Extract text, detect objects, and generate captions from images""",
        },
        {
            "layout": "content",
            "title": "What is Azure AI Vision?",
            "content": """Cloud-based AI service that analyzes images and extracts valuable information

Part of Microsoft Foundry Tools (Azure AI Services)

Key Capabilities:
- Read and extract text from images (OCR)
- Detect and identify objects in photos
- Generate human-like image descriptions
- Analyze image content and metadata

Pricing:
- Free tier: 20 calls/min, 5,000 calls/month
- Perfect for learning and prototyping!

learn.microsoft.com/azure/ai-services/computer-vision/""",
        },
        {
            "layout": "content_alt",
            "title": "Key Features",
            "content": """OCR (Read):
- Extract printed & handwritten text
- Scan receipts, documents, signs

Image Captions:
- Generate natural language descriptions
- "A dog playing in the park"

Object Detection:
- Identify objects with bounding boxes
- Find all cars in a parking lot

Smart Tags:
- Label images with keywords
- ["outdoor", "nature", "sunny"]

Dense Captions:
- Multiple captions for different image regions""",
        },
        {
            "layout": "content",
            "title": "Real-World Use Cases",
            "content": """Accessibility:
- Screen readers for visually impaired users
- Auto-generate alt-text for websites

Retail & E-commerce:
- Product image tagging
- Visual search ("find similar items")

Document Processing:
- Digitize paper documents
- Extract data from forms and receipts

Security:
- Object detection in security feeds
- License plate recognition

Social Media:
- Content moderation
- Auto-tagging photos""",
        },
        {
            "layout": "content_alt2",
            "title": "What We'll Build",
            "content": """A Streamlit web app that:

- Accepts image upload or URL
- Sends image to Azure AI Vision API
- Displays AI-powered results

Features we'll implement:
- Generated caption describing the image
- Detected objects with bounding boxes
- Extracted text (OCR)
- Smart tags and keywords

Tech Stack:
- Python 3.10+
- Streamlit (Web UI)
- Azure AI Vision REST API
- About 50 lines of code!""",
        },
        {
            "layout": "demo",
            "title": "Demo Time!",
            "subtitle": "Let's build it live!",
        },
        {
            "layout": "content",
            "title": "Resources & Next Steps",
            "content": """Official Documentation:
- learn.microsoft.com/azure/ai-services/computer-vision/

Microsoft Learn Path:
- learn.microsoft.com/training/paths/create-computer-vision-solutions-azure-ai/

Certifications:
- AI-102: Azure AI Engineer Associate
- AI-900: Azure AI Fundamentals

Code from Today:
- github.com/jainal09/azure-ai-vision-workshop

Connect with Me:
- linkedin.com/in/jainal-gosaliya""",
        },
        {
            "layout": "section",
            "title": "Q&A",
        },
        {
            "layout": "closing",
            "title": "Thank You!",
        },
    ]

    # Add slides
    for i, slide_data in enumerate(slides_content):
        print(f"Adding slide {i + 2}: {slide_data.get('title', 'Untitled')} ({slide_data['layout']})")
        add_slide(prs, slide_data)

    # Save
    prs.save(OUTPUT_PATH)
    print(f"\nPresentation saved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    generate_presentation()
